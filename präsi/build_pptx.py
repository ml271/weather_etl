"""
Build Weather ETL Capstone Presentation as .pptx
Run: python build_pptx.py
Output: weather_etl_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x14, 0x24)   # dark navy
SURFACE  = RGBColor(0x16, 0x21, 0x38)   # slightly lighter
SURFACE2 = RGBColor(0x1e, 0x2d, 0x4a)   # card bg
ACCENT   = RGBColor(0x4f, 0xc3, 0xf7)   # sky blue
ACCENT2  = RGBColor(0x81, 0xc7, 0x84)   # green
ACCENT3  = RGBColor(0xff, 0xb7, 0x4d)   # orange
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x90, 0xa4, 0xae)
BORDER   = RGBColor(0x2a, 0x3f, 0x5f)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Full dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill_color=SURFACE2, border_color=BORDER, border_pt=1.0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, x, y, w, h, size=9, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    """Small uppercase label"""
    return txbox(slide, text.upper(), x, y, w, h, size=size, color=color,
                 bold=bold, align=align)

def card(slide, title, body, x, y, w, h,
         title_color=WHITE, body_color=MUTED, fill=SURFACE2, border=BORDER):
    rect(slide, x, y, w, h, fill_color=fill, border_color=border)
    txbox(slide, title, x + Inches(0.12), y + Inches(0.1), w - Inches(0.24), Inches(0.35),
          size=13, bold=True, color=title_color)
    txbox(slide, body,  x + Inches(0.12), y + Inches(0.38), w - Inches(0.24), h - Inches(0.5),
          size=10, color=body_color)

def section_label(slide, text, x=Inches(0.5), y=Inches(0.25)):
    txbox(slide, text.upper(), x, y, Inches(6), Inches(0.3),
          size=9, color=ACCENT, bold=False)

def slide_title(slide, text, x=Inches(0.5), y=Inches(0.5)):
    txbox(slide, text, x, y, Inches(12.3), Inches(0.65),
          size=32, bold=True, color=WHITE)

def divider(slide, y, x=Inches(0.5), w=None):
    if w is None:
        w = Inches(12.33)
    ln = slide.shapes.add_shape(1, x, y, w, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BORDER
    ln.line.fill.background()

def bullet_box(slide, items, x, y, w, h, title=None, title_color=ACCENT):
    if title:
        txbox(slide, title, x, y, w, Inches(0.3), size=12, bold=True, color=title_color)
        y += Inches(0.32)
        h -= Inches(0.32)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED

def flow_boxes(slide, steps, y, arrow_color=ACCENT):
    """Horizontal flow: [(icon, label), ...]"""
    n = len(steps)
    arrow_w = Inches(0.4)
    box_w   = Inches(1.2)
    total_w = n * box_w + (n - 1) * arrow_w
    start_x = (W - total_w) / 2

    for i, (icon, lbl) in enumerate(steps):
        bx = start_x + i * (box_w + arrow_w)
        rect(slide, bx, y, box_w, Inches(0.7), fill_color=SURFACE2, border_color=BORDER)
        txbox(slide, icon, bx, y + Inches(0.03), box_w, Inches(0.32),
              size=16, align=PP_ALIGN.CENTER)
        txbox(slide, lbl, bx, y + Inches(0.36), box_w, Inches(0.28),
              size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = bx + box_w
            txbox(slide, "→", ax, y + Inches(0.2), arrow_w, Inches(0.3),
                  size=14, color=arrow_color, align=PP_ALIGN.CENTER)

def schema_table(slide, headers, rows, x, y, w, row_h=Inches(0.28)):
    col_w = w / len(headers)
    # Header row
    for ci, h in enumerate(headers):
        cx = x + ci * col_w
        rect(slide, cx, y, col_w, row_h,
             fill_color=RGBColor(0x1a, 0x36, 0x52), border_color=BORDER)
        txbox(slide, h, cx + Inches(0.06), y + Inches(0.05), col_w - Inches(0.1), row_h,
              size=9, bold=True, color=ACCENT)
    # Data rows
    for ri, row in enumerate(rows):
        ry = y + (ri + 1) * row_h
        for ci, cell in enumerate(row):
            cx = x + ci * col_w
            fill = SURFACE2 if ri % 2 == 0 else RGBColor(0x18, 0x28, 0x40)
            rect(slide, cx, ry, col_w, row_h, fill_color=fill, border_color=BORDER)
            txbox(slide, cell, cx + Inches(0.06), ry + Inches(0.04),
                  col_w - Inches(0.1), row_h, size=9, color=MUTED)

def add_image(slide, path, x, y, w, h):
    full_path = os.path.join(SCRIPT_DIR, path)
    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, x, y, w, h)
    else:
        # placeholder rectangle
        rect(slide, x, y, w, h, fill_color=SURFACE2, border_color=ACCENT)
        txbox(slide, f"📷 {path}", x, y + h/2 - Inches(0.2), w, Inches(0.4),
              size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()

# Gradient-style accent bar top
r = rect(s, 0, 0, W, Inches(0.08), fill_color=ACCENT, border_color=None)

# Main title
txbox(s, "Weather ETL System", Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
      size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
txbox(s, "A real-time weather data pipeline with automated alerts & dashboarding",
      Inches(1.5), Inches(2.8), Inches(10.33), Inches(0.5),
      size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Tech pills row
pills = ["Apache Airflow", "PostgreSQL", "FastAPI", "Docker", "Nginx", "Open-Meteo API"]
pill_w = Inches(1.7)
pill_gap = Inches(0.15)
total = len(pills) * pill_w + (len(pills)-1) * pill_gap
px = (W - total) / 2
for i, p in enumerate(pills):
    color = ACCENT2 if p == "Open-Meteo API" else ACCENT
    bcolor = ACCENT2 if p == "Open-Meteo API" else ACCENT
    bg = RGBColor(0x1a, 0x36, 0x52) if p != "Open-Meteo API" else RGBColor(0x1a, 0x36, 0x2a)
    rx = px + i * (pill_w + pill_gap)
    r2 = rect(s, rx, Inches(3.5), pill_w, Inches(0.38), fill_color=bg, border_color=bcolor)
    txbox(s, p, rx, Inches(3.55), pill_w, Inches(0.3),
          size=11, color=color, align=PP_ALIGN.CENTER, bold=True)

# Author box
rect(s, Inches(4.0), Inches(4.2), Inches(5.33), Inches(1.1),
     fill_color=SURFACE2, border_color=ACCENT)
txbox(s, "Marvin Lorff", Inches(4.0), Inches(4.35), Inches(5.33), Inches(0.45),
      size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "neuefische · Data Engineering Bootcamp · Capstone Project",
      Inches(4.0), Inches(4.75), Inches(5.33), Inches(0.35),
      size=11, color=MUTED, align=PP_ALIGN.CENTER)

txbox(s, "March 2026", Inches(0), Inches(6.9), W, Inches(0.3),
      size=10, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill_color=ACCENT2, border_color=None)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Overview")
slide_title(s, "What is Weather ETL?")
divider(s, Inches(1.22))

txbox(s, "An end-to-end data engineering project that automatically fetches, processes, and stores "
         "weather forecasts — and notifies users when their custom alert conditions are met.",
      Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.55), size=13, color=MUTED)

# 4 stat cards
stats = [
    ("🌍", "~75 Cities", "Germany, Austria & Switzerland"),
    ("⏱️", "Hourly Runs", "Automated pipeline schedule"),
    ("📊", "7-Day Forecast", "Daily + 168h hourly data"),
    ("🔔", "Email Alerts", "Triggered notifications"),
]
cw = Inches(2.8)
for i, (icon, title, body) in enumerate(stats):
    cx = Inches(0.5) + i * (cw + Inches(0.22))
    rect(s, cx, Inches(2.05), cw, Inches(1.1), fill_color=SURFACE2, border_color=BORDER)
    txbox(s, icon, cx, Inches(2.1), cw, Inches(0.4), size=20, align=PP_ALIGN.CENTER)
    txbox(s, title, cx, Inches(2.5), cw, Inches(0.32),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, body, cx, Inches(2.8), cw, Inches(0.28),
          size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Flow diagram
txbox(s, "Data Flow", Inches(0.5), Inches(3.35), Inches(3), Inches(0.3),
      size=10, color=MUTED)
flow_boxes(s, [
    ("🌐", "Open-Meteo API"),
    ("⚙️", "Airflow ETL"),
    ("🗄️", "PostgreSQL"),
    ("🚀", "FastAPI"),
    ("🖥️", "Dashboard"),
], Inches(3.75))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Technology Stack
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Stack")
slide_title(s, "Technology Stack")
divider(s, Inches(1.22))

tech = [
    ("⚙️ Apache Airflow",  "Pipeline orchestration · DAG scheduling · XCom data passing · Retry logic"),
    ("🗄️ PostgreSQL 15",   "Primary datastore · 6 tables · JSONB conditions · UPSERT pattern"),
    ("🚀 FastAPI",          "REST API · JWT auth · 18 endpoints · SQLAlchemy ORM"),
    ("🌐 Nginx + JS",       "Static frontend · Station search · Chart.js dashboard"),
    ("🐳 Docker Compose",   "6 services · Isolated containers · Single-command deployment"),
    ("🌤️ Open-Meteo API",  "Free · No API key required · 7-day forecast · Daily + hourly"),
]
cw = Inches(4.0)
ch = Inches(1.3)
gap = Inches(0.22)
for i, (title, body) in enumerate(tech):
    col = i % 3
    row = i // 3
    cx = Inches(0.5) + col * (cw + gap)
    cy = Inches(1.4) + row * (ch + Inches(0.18))
    rect(s, cx, cy, cw, ch, fill_color=SURFACE2, border_color=BORDER)
    txbox(s, title, cx + Inches(0.15), cy + Inches(0.12), cw - Inches(0.3), Inches(0.38),
          size=14, bold=True, color=WHITE)
    txbox(s, body,  cx + Inches(0.15), cy + Inches(0.5),  cw - Inches(0.3), Inches(0.7),
          size=11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ETL Pipeline Overview
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "ETL Pipeline")
slide_title(s, "The Pipeline at a Glance")
divider(s, Inches(1.22))

# DAG 1 box
rect(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(1.35),
     fill_color=SURFACE2, border_color=ACCENT)
txbox(s, "DAG 1 — weather_etl  (every 2h)",
      Inches(0.65), Inches(1.42), Inches(5.6), Inches(0.3),
      size=11, bold=True, color=ACCENT)
dag1 = [("📡","Extract"), ("🔄","Transform"), ("💾","Load")]
for i, (icon, lbl) in enumerate(dag1):
    bx = Inches(0.7) + i * Inches(1.7)
    rect(s, bx, Inches(1.78), Inches(1.3), Inches(0.65),
         fill_color=RGBColor(0x12,0x22,0x38), border_color=BORDER)
    txbox(s, icon, bx, Inches(1.8), Inches(1.3), Inches(0.3),
          size=14, align=PP_ALIGN.CENTER)
    txbox(s, lbl, bx, Inches(2.1), Inches(1.3), Inches(0.25),
          size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2:
        txbox(s, "→", bx + Inches(1.3), Inches(1.93), Inches(0.4), Inches(0.3),
              size=14, color=ACCENT, align=PP_ALIGN.CENTER)

# DAG 2 box
rect(s, Inches(6.93), Inches(1.35), Inches(5.9), Inches(1.35),
     fill_color=SURFACE2, border_color=ACCENT2)
txbox(s, "DAG 2 — check_warnings  (every 1h)",
      Inches(7.08), Inches(1.42), Inches(5.6), Inches(0.3),
      size=11, bold=True, color=ACCENT2)
dag2 = [("🗄️","Read DB"), ("⚖️","Evaluate"), ("📧","Notify")]
for i, (icon, lbl) in enumerate(dag2):
    bx = Inches(7.15) + i * Inches(1.7)
    rect(s, bx, Inches(1.78), Inches(1.3), Inches(0.65),
         fill_color=RGBColor(0x12,0x28,0x22), border_color=RGBColor(0x2a,0x4f,0x3a))
    txbox(s, icon, bx, Inches(1.8), Inches(1.3), Inches(0.3),
          size=14, align=PP_ALIGN.CENTER)
    txbox(s, lbl, bx, Inches(2.1), Inches(1.3), Inches(0.25),
          size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2:
        txbox(s, "→", bx + Inches(1.3), Inches(1.93), Inches(0.4), Inches(0.3),
              size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

# E / T / L cards side by side
etl = [
    ("E", "Extract", ACCENT,
     "GET api.open-meteo.com with lat/lon. Returns 7-day daily (10 vars) + "
     "168h hourly (17 vars incl. soil data). Raw JSON pushed to Airflow XCom."),
    ("T", "Transform", ACCENT3,
     "API arrays unpacked into flat row dicts. WMO codes mapped to descriptions. "
     "Alert rules from YAML evaluated with AND logic against daily records."),
    ("L", "Load", ACCENT2,
     "ON CONFLICT DO UPDATE — no duplicate rows. Previous alerts deactivated; "
     "fresh alerts inserted per run. Fully idempotent across all pipeline runs."),
]
ew = Inches(4.0)
for i, (letter, title, color, body) in enumerate(etl):
    ex = Inches(0.5) + i * (ew + Inches(0.22))
    ey = Inches(2.9)
    rect(s, ex, ey, ew, Inches(1.95), fill_color=SURFACE2, border_color=color)
    txbox(s, letter, ex + Inches(0.15), ey + Inches(0.1), Inches(0.55), Inches(0.6),
          size=36, bold=True, color=color)
    txbox(s, title, ex + Inches(0.72), ey + Inches(0.22), ew - Inches(0.9), Inches(0.35),
          size=14, bold=True, color=WHITE)
    txbox(s, body, ex + Inches(0.15), ey + Inches(0.72), ew - Inches(0.3), Inches(1.1),
          size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Extract Detail
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "ETL — Extract")
slide_title(s, "Extracting Weather Data")
divider(s, Inches(1.22))

# Left column
rect(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(5.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Daily Variables — 7 days", Inches(0.65), Inches(1.45), Inches(5.6), Inches(0.35),
      size=13, bold=True, color=ACCENT)
daily = [
    "Temperature max / min  (°C)",
    "Precipitation sum  (mm)",
    "Snowfall sum  (cm)",
    "Wind speed max  (km/h)",
    "Wind gusts max  (km/h)",
    "UV Index max",
    "WMO weather code",
    "Sunrise / Sunset times",
]
for i, item in enumerate(daily):
    txbox(s, "• " + item, Inches(0.7), Inches(1.85) + i * Inches(0.55),
          Inches(5.5), Inches(0.45), size=12, color=MUTED)

# Right column
rect(s, Inches(6.93), Inches(1.35), Inches(5.9), Inches(5.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Hourly Variables — 168 hours", Inches(7.08), Inches(1.45), Inches(5.6), Inches(0.35),
      size=13, bold=True, color=ACCENT)
hourly = [
    "Temperature & apparent temperature  (°C)",
    "Precipitation / Rain / Snowfall",
    "Wind speed & direction",
    "Relative humidity  (%)",
    "Sunshine duration  (sec/h)",
    "Soil temperature: 0 cm / 6 cm / 18 cm  (°C)",
    "Soil moisture: 0–1 / 1–3 / 3–9 cm  (m³/m³)",
]
for i, item in enumerate(hourly):
    txbox(s, "• " + item, Inches(7.12), Inches(1.85) + i * Inches(0.62),
          Inches(5.6), Inches(0.5), size=12, color=MUTED)

# Footer info
rect(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.35),
     fill_color=RGBColor(0x12,0x22,0x38), border_color=BORDER)
txbox(s, "XCom payload:  city · latitude · longitude · fetched_at (UTC) · full Open-Meteo JSON  |  "
         "API: Free tier · No API key required · Timezone: Europe/Berlin",
      Inches(0.65), Inches(6.98), Inches(12.0), Inches(0.28),
      size=9, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Transform Detail
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "ETL — Transform")
slide_title(s, "Transforming the Raw Data")
divider(s, Inches(1.22))

steps = [
    ("1", "Unpack API Arrays → Flat Row Dicts",
     "Open-Meteo returns columnar arrays (e.g. temperature_2m_max: [12.1, 15.3, …]).\n"
     "Transform zips all arrays by index into individual row dicts — one dict per calendar day or hour."),
    ("2", "WMO Code Enrichment",
     "Each row's numeric WMO weather code (e.g. 63) is translated to a human-readable\n"
     "string (\"Moderate Rain\") using a built-in lookup table covering all 30 WMO codes."),
    ("3", "Alert Rule Evaluation  (YAML Config)",
     "Alert thresholds loaded from alerts_config.yaml — not hardcoded in Python.\n"
     "Each rule defines: name, severity, conditions (AND logic), message text.\n"
     "Triggered alerts are included in the XCom payload alongside daily/hourly records."),
]
for i, (num, title, body) in enumerate(steps):
    ey = Inches(1.4) + i * Inches(1.72)
    rect(s, Inches(0.5), ey, Inches(12.33), Inches(1.55),
         fill_color=SURFACE2, border_color=BORDER)
    txbox(s, num, Inches(0.6), ey + Inches(0.1), Inches(0.55), Inches(0.8),
          size=40, bold=True, color=ACCENT)
    txbox(s, title, Inches(1.25), ey + Inches(0.15), Inches(11.4), Inches(0.38),
          size=14, bold=True, color=WHITE)
    txbox(s, body, Inches(1.25), ey + Inches(0.58), Inches(11.4), Inches(0.9),
          size=11, color=MUTED)

# Alert examples
rect(s, Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.55),
     fill_color=RGBColor(0x1a, 0x28, 0x1a), border_color=ACCENT2)
examples = "Example rules:   temperature_max > 35  →  danger   |   precipitation_sum > 20  →  warning   |   wind_speed_max > 60  →  danger   |   snowfall_sum > 10  →  info"
txbox(s, examples, Inches(0.65), Inches(6.72), Inches(12.0), Inches(0.35),
      size=10, color=ACCENT2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Load Detail
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "ETL — Load")
slide_title(s, "Loading Data into PostgreSQL")
divider(s, Inches(1.22))

# Left
rect(s, Inches(0.5), Inches(1.35), Inches(5.9), Inches(2.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "UPSERT Pattern", Inches(0.65), Inches(1.45), Inches(5.6), Inches(0.35),
      size=14, bold=True, color=ACCENT)
upsert = [
    "Unique constraint: (city, forecast_date)  for daily",
    "Unique constraint: (city, forecast_time)  for hourly",
    "On conflict → update all columns in-place",
    "No duplicate rows accumulate across runs",
]
for i, item in enumerate(upsert):
    txbox(s, "• " + item, Inches(0.7), Inches(1.88) + i * Inches(0.48),
          Inches(5.5), Inches(0.4), size=12, color=MUTED)

# Right
rect(s, Inches(6.93), Inches(1.35), Inches(5.9), Inches(2.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Alert Lifecycle", Inches(7.08), Inches(1.45), Inches(5.6), Inches(0.35),
      size=14, bold=True, color=ACCENT3)
alert_l = [
    "Before insert: all previous alerts for city",
    "  set  is_active = false",
    "Fresh alerts from current run inserted as",
    "  is_active = true",
    "Dashboard always shows current forecast only",
]
for i, item in enumerate(alert_l):
    txbox(s, item, Inches(7.12), Inches(1.88) + i * Inches(0.42),
          Inches(5.6), Inches(0.38), size=12, color=MUTED)

# Bottom summary box
rect(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(1.1),
     fill_color=RGBColor(0x12,0x22,0x38), border_color=ACCENT)
txbox(s, "What gets written per pipeline run:",
      Inches(0.65), Inches(4.18), Inches(4), Inches(0.3), size=11, bold=True, color=ACCENT)
write_items = [
    ("7 rows", "weather_daily"),
    ("168 rows", "weather_hourly"),
    ("n rows (if triggered)", "weather_alerts"),
]
for i, (count, table) in enumerate(write_items):
    bx = Inches(0.5) + i * Inches(4.11)
    txbox(s, count, bx, Inches(4.52), Inches(4.0), Inches(0.28),
          size=10, color=MUTED, align=PP_ALIGN.CENTER)
    txbox(s, table, bx, Inches(4.77), Inches(4.0), Inches(0.3),
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PostgreSQL Schema Overview
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Database")
slide_title(s, "PostgreSQL Schema — 6 Tables")
divider(s, Inches(1.22))

tables = [
    ("📅", "weather_daily",    "13 columns", ACCENT,  "Daily forecast per city · UPSERT on (city, date)"),
    ("🕐", "weather_hourly",   "21 columns", ACCENT,  "Hourly forecast + soil data · UPSERT on (city, time)"),
    ("🔔", "weather_alerts",   "10 columns", ACCENT,  "System-generated alerts · JSONB trigger snapshot"),
    ("📍", "stations",         " 6 columns", ACCENT2, "~75 cities pre-seeded · lat/lon for API calls"),
    ("👤", "users",            " 6 columns", ACCENT2, "bcrypt-hashed passwords · email for notifications"),
    ("⚠️", "warnings",         "10 columns", ACCENT2, "User rules · JSONB conditions + validity · active flag"),
]
tw = Inches(3.95)
th = Inches(1.5)
gap = Inches(0.22)
for i, (icon, name, cols, color, desc) in enumerate(tables):
    col = i % 3
    row = i // 3
    tx = Inches(0.5) + col * (tw + gap)
    ty = Inches(1.4) + row * (th + Inches(0.2))
    border_c = color if color == ACCENT2 else BORDER
    rect(s, tx, ty, tw, th, fill_color=SURFACE2, border_color=border_c)
    txbox(s, icon, tx + Inches(0.12), ty + Inches(0.1), Inches(0.55), Inches(0.5),
          size=22)
    txbox(s, name, tx + Inches(0.72), ty + Inches(0.13), tw - Inches(0.85), Inches(0.35),
          size=13, bold=True, color=WHITE)
    txbox(s, cols, tx + Inches(0.72), ty + Inches(0.47), tw - Inches(0.85), Inches(0.25),
          size=10, color=color)
    txbox(s, desc, tx + Inches(0.12), ty + Inches(0.82), tw - Inches(0.24), Inches(0.6),
          size=10, color=MUTED)

# Legend
txbox(s, "ETL tables    |    Application tables",
      Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.3),
      size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Schema Detail: Core ETL Tables
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Database — Schema Detail")
slide_title(s, "Core ETL Tables")
divider(s, Inches(1.22))

# weather_daily
txbox(s, "weather_daily", Inches(0.5), Inches(1.32), Inches(6), Inches(0.35),
      size=14, bold=True, color=ACCENT)
schema_table(s,
    ["Column", "Type", "Unit"],
    [
        ["id  (PK)", "INTEGER", "—"],
        ["city", "VARCHAR(100)", "—"],
        ["forecast_date", "DATE", "—"],
        ["temperature_max", "DECIMAL(5,2)", "°C"],
        ["temperature_min", "DECIMAL(5,2)", "°C"],
        ["precipitation_sum", "DECIMAL(6,2)", "mm"],
        ["snowfall_sum", "DECIMAL(6,2)", "cm"],
        ["wind_speed_max", "DECIMAL(6,2)", "km/h"],
        ["wind_gusts_max", "DECIMAL(6,2)", "km/h"],
        ["uv_index_max", "DECIMAL(4,2)", "—"],
        ["weather_code", "INTEGER", "WMO"],
        ["sunrise / sunset", "TIMESTAMPTZ", "UTC"],
        ["created_at", "TIMESTAMPTZ", "UTC"],
    ],
    Inches(0.5), Inches(1.7), Inches(6.0),
    row_h=Inches(0.33)
)

# weather_hourly (selection)
txbox(s, "weather_hourly  (key columns)", Inches(6.93), Inches(1.32), Inches(6), Inches(0.35),
      size=14, bold=True, color=ACCENT)
schema_table(s,
    ["Column", "Type", "Unit"],
    [
        ["temperature", "DECIMAL(5,2)", "°C"],
        ["feels_like", "DECIMAL(5,2)", "°C"],
        ["precipitation", "DECIMAL(5,2)", "mm"],
        ["wind_speed", "DECIMAL(6,2)", "km/h"],
        ["wind_direction", "INTEGER", "degrees"],
        ["humidity", "INTEGER", "%"],
        ["sunshine_duration", "DECIMAL(6,2)", "sec"],
        ["soil_temperature_0cm", "DECIMAL(5,2)", "°C"],
        ["soil_temperature_6cm", "DECIMAL(5,2)", "°C"],
        ["soil_temperature_18cm", "DECIMAL(5,2)", "°C"],
        ["soil_moisture_0_1cm", "DECIMAL(8,6)", "m³/m³"],
        ["soil_moisture_1_3cm", "DECIMAL(8,6)", "m³/m³"],
        ["soil_moisture_3_9cm", "DECIMAL(8,6)", "m³/m³"],
    ],
    Inches(6.93), Inches(1.7), Inches(5.9),
    row_h=Inches(0.33)
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Schema Detail: Application Tables
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Database — Schema Detail")
slide_title(s, "Application Tables")
divider(s, Inches(1.22))

# users
txbox(s, "users", Inches(0.5), Inches(1.32), Inches(6), Inches(0.35),
      size=14, bold=True, color=ACCENT2)
schema_table(s,
    ["Column", "Type / Notes"],
    [
        ["id  (PK)",         "INTEGER"],
        ["email",            "VARCHAR(255)  UNIQUE"],
        ["username",         "VARCHAR(100)  UNIQUE"],
        ["hashed_password",  "VARCHAR(255)  — bcrypt"],
        ["is_active",        "BOOLEAN"],
        ["created_at",       "TIMESTAMPTZ"],
    ],
    Inches(0.5), Inches(1.7), Inches(5.9),
    row_h=Inches(0.38)
)

# warnings
txbox(s, "warnings  (user-defined)", Inches(6.93), Inches(1.32), Inches(6), Inches(0.35),
      size=14, bold=True, color=ACCENT2)
schema_table(s,
    ["Column", "Notes"],
    [
        ["user_id  (FK)", "→ users.id  CASCADE DELETE"],
        ["station_id  (FK)", "→ stations.id  SET NULL"],
        ["city", "Target city name"],
        ["name", "User-chosen label"],
        ["conditions", "JSONB  array of rules"],
        ["validity", "JSONB  — date_range / weekdays / months"],
        ["active", "Enable / disable toggle"],
    ],
    Inches(6.93), Inches(1.7), Inches(5.9),
    row_h=Inches(0.38)
)

# weather_alerts note
rect(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(1.85),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "weather_alerts", Inches(0.65), Inches(5.4), Inches(5), Inches(0.32),
      size=13, bold=True, color=ACCENT)
alert_cols = [
    ("alert_name", "From YAML config"),
    ("severity", "info / warning / danger"),
    ("message", "Human-readable description"),
    ("condition_met", "JSONB — value + threshold snapshot"),
    ("is_active", "Only latest pipeline run = true"),
]
for i, (col, note) in enumerate(alert_cols):
    row_x = Inches(0.65) + (i % 3) * Inches(4.0)
    row_y = Inches(5.78) + (i // 3) * Inches(0.42)
    txbox(s, f"• {col}:", row_x, row_y, Inches(1.6), Inches(0.35),
          size=10, bold=True, color=MUTED)
    txbox(s, note, row_x + Inches(1.55), row_y, Inches(2.35), Inches(0.35),
          size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — User Warning Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Feature")
slide_title(s, "User Warning Pipeline")
divider(s, Inches(1.22))

txbox(s, "Users configure custom alert rules — the second Airflow DAG checks them hourly and sends email notifications.",
      Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.4), size=12, color=MUTED)

flow_boxes(s, [
    ("🖱️", "User creates\nwarning in UI"),
    ("🗄️", "Stored as JSONB\nin warnings table"),
    ("⚙️", "Airflow DAG\nruns every hour"),
    ("⚖️", "Conditions checked\nvs. weather_daily"),
    ("📧", "Email sent\nvia SMTP"),
], Inches(1.85))

# Two columns below
rect(s, Inches(0.5), Inches(3.0), Inches(5.9), Inches(3.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Condition Types", Inches(0.65), Inches(3.1), Inches(5.6), Inches(0.35),
      size=13, bold=True, color=ACCENT)
cond_items = [
    "Parameters: temp max/min, precipitation, snowfall,\n  wind, gusts, UV index",
    "Comparators: > >= < <= ==",
    "Logic: AND across all conditions in a rule",
    "",
    "Validity Modes:",
    "  date_range  — from / to specific dates",
    "  weekdays     — e.g. Saturday & Sunday only",
    "  months        — e.g. January & February only",
]
for i, item in enumerate(cond_items):
    txbox(s, item, Inches(0.7), Inches(3.52) + i * Inches(0.36),
          Inches(5.5), Inches(0.32), size=11, color=MUTED)

rect(s, Inches(6.93), Inches(3.0), Inches(5.9), Inches(3.5),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Pre-built Templates", Inches(7.08), Inches(3.1), Inches(5.6), Inches(0.35),
      size=13, bold=True, color=ACCENT3)
templates = [
    "🌡️  Heatwave:    temperature_max > 32°C",
    "❄️   Frost:          temperature_min < 0°C",
    "🌧️  Heavy Rain:  precipitation_sum > 20 mm",
    "💨  Storm:          wind_speed_max > 60 km/h",
    "🌨️  Snowfall:      snowfall_sum > 10 cm",
]
for i, t in enumerate(templates):
    txbox(s, t, Inches(7.12), Inches(3.52) + i * Inches(0.6),
          Inches(5.6), Inches(0.5), size=12, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Docker Network Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Infrastructure")
slide_title(s, "Docker Network Architecture")
divider(s, Inches(1.22))

txbox(s, "All 6 services run in a single Docker Compose stack — one command to deploy everything.",
      Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.35), size=12, color=MUTED)

services = [
    ("🌐", "Nginx Frontend",    ":80",          ACCENT,  "Static HTML/JS · Station search · Dashboard"),
    ("🚀", "FastAPI Backend",   ":8000",         ACCENT,  "REST API · JWT Auth · SQLAlchemy ORM"),
    ("🗄️", "PostgreSQL 15",     ":5432",         ACCENT,  "Primary datastore · Persistent volume"),
    ("⚙️", "Airflow Webserver", ":8080",         ACCENT2, "DAG monitoring · Run history"),
    ("🕐", "Airflow Scheduler", "internal",      ACCENT2, "DAG execution · Retry policy"),
    ("📬", "Mailhog / SMTP",    ":1025 / :8025", ACCENT2, "SMTP catch-all · Email preview UI"),
]
sw = Inches(4.0)
sh = Inches(1.3)
gap = Inches(0.22)
for i, (icon, name, port, color, desc) in enumerate(services):
    col = i % 3
    row = i // 3
    sx = Inches(0.5) + col * (sw + gap)
    sy = Inches(1.75) + row * (sh + Inches(0.2))
    rect(s, sx, sy, sw, sh, fill_color=SURFACE2, border_color=color)
    txbox(s, icon, sx + Inches(0.12), sy + Inches(0.1), Inches(0.5), Inches(0.45), size=20)
    txbox(s, name, sx + Inches(0.7), sy + Inches(0.12), sw - Inches(0.85), Inches(0.32),
          size=13, bold=True, color=WHITE)
    txbox(s, port, sx + Inches(0.7), sy + Inches(0.44), sw - Inches(0.85), Inches(0.25),
          size=10, color=color)
    txbox(s, desc, sx + Inches(0.12), sy + Inches(0.72), sw - Inches(0.24), Inches(0.5),
          size=10, color=MUTED)

rect(s, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.38),
     fill_color=RGBColor(0x12,0x22,0x38), border_color=BORDER)
txbox(s, "All services share one internal Docker network  ·  Secrets via .env  ·  Backend runs as non-root user  ·  Data persisted in named volumes",
      Inches(0.65), Inches(6.96), Inches(12.0), Inches(0.28), size=9, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Dashboard Screenshots
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Demo")
slide_title(s, "Dashboard & Features")
divider(s, Inches(1.22))

screenshots = [
    ("Screenshot 2026-03-19 091336.png", "Station Search"),
    ("Screenshot 2026-03-19 092656.png", "Weather Dashboard"),
    ("Screenshot 2026-03-19 092752.png", "Warning Builder"),
    ("Screenshot 2026-03-19 092809.png", "Active Warnings"),
]
iw = Inches(6.0)
ih = Inches(2.7)
for i, (fname, caption) in enumerate(screenshots):
    col = i % 2
    row = i // 2
    ix = Inches(0.5) + col * (iw + Inches(0.33))
    iy = Inches(1.4) + row * (ih + Inches(0.55))
    add_image(s, fname, ix, iy, iw, ih)
    txbox(s, caption, ix, iy + ih + Inches(0.05), iw, Inches(0.3),
          size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Key Numbers
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
section_label(s, "Summary")
slide_title(s, "By the Numbers")
divider(s, Inches(1.22))

numbers = [
    ("~75",  "Pre-seeded\nstations"),
    ("6",    "PostgreSQL\ntables"),
    ("175",  "Rows per\npipeline run"),
    ("2",    "Airflow\nDAGs"),
    ("18",   "API\nendpoints"),
    ("6",    "Docker\nservices"),
]
nw = Inches(2.0)
nx_start = (W - (6 * nw + 5 * Inches(0.12))) / 2
for i, (num, lbl) in enumerate(numbers):
    nx = nx_start + i * (nw + Inches(0.12))
    rect(s, nx, Inches(1.5), nw, Inches(1.8), fill_color=SURFACE2, border_color=BORDER)
    txbox(s, num, nx, Inches(1.6), nw, Inches(0.85),
          size=46, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txbox(s, lbl, nx, Inches(2.42), nw, Inches(0.7),
          size=11, color=MUTED, align=PP_ALIGN.CENTER)

highlights = [
    "✅  OWASP security hardening — JWT, bcrypt, non-root container, XSS prevention",
    "✅  UPSERT pattern — fully idempotent pipeline runs",
    "✅  JSONB conditions — flexible warning rules without schema migrations",
    "✅  YAML alert config — threshold changes without code redeploy",
]
rect(s, Inches(0.5), Inches(3.55), Inches(12.33), Inches(2.5),
     fill_color=SURFACE2, border_color=BORDER)
for i, h in enumerate(highlights):
    txbox(s, h, Inches(0.7), Inches(3.7) + i * Inches(0.55),
          Inches(12.0), Inches(0.45), size=13, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()

rect(s, 0, 0, W, Inches(0.08), fill_color=ACCENT, border_color=None)

txbox(s, "Thank You", Inches(1), Inches(1.6), Inches(11.33), Inches(1.1),
      size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Weather ETL — Capstone Project",
      Inches(1), Inches(2.75), Inches(11.33), Inches(0.5),
      size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Two cards
cw2 = Inches(4.5)
for col, (icon, title, body) in enumerate([
    ("💻", "GitHub", "github.com/ml271/weather_etl"),
    ("📚", "Stack", "Airflow · PostgreSQL · FastAPI · Docker · Nginx"),
]):
    cx = Inches(2.16) + col * (cw2 + Inches(0.5))
    rect(s, cx, Inches(3.5), cw2, Inches(1.2), fill_color=SURFACE2, border_color=ACCENT)
    txbox(s, icon + "  " + title, cx, Inches(3.62), cw2, Inches(0.42),
          size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, body, cx, Inches(4.05), cw2, Inches(0.4),
          size=11, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(5.0), Inches(6.33), Inches(1.1),
     fill_color=SURFACE2, border_color=BORDER)
txbox(s, "Marvin Lorff", Inches(3.5), Inches(5.12), Inches(6.33), Inches(0.42),
      size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "neuefische · Data Engineering Bootcamp · Capstone Project · March 2026",
      Inches(3.5), Inches(5.52), Inches(6.33), Inches(0.35),
      size=11, color=MUTED, align=PP_ALIGN.CENTER)

txbox(s, "Questions?", Inches(0), Inches(6.5), W, Inches(0.5),
      size=14, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill_color=ACCENT2, border_color=None)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(SCRIPT_DIR, "weather_etl_presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
